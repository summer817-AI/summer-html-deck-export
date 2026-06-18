import argparse
import json
import math
import shutil
import subprocess
import sys
import time
from pathlib import Path

from bs4 import BeautifulSoup
from PIL import Image, ImageEnhance
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pypdf import PdfReader


DEFAULT_CHROME_PATHS = [
    r"C:\Program Files\Google\Chrome\Application\chrome.exe",
    r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
    r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
]

DARK_STOPS = [
    (0.00, (2, 7, 16, int(0.90 * 255))),
    (0.52, (7, 23, 43, int(0.64 * 255))),
    (1.00, (2, 7, 16, int(0.24 * 255))),
]
LIGHT_STOPS = [
    (0.00, (244, 248, 250, int(0.94 * 255))),
    (0.54, (244, 248, 250, int(0.80 * 255))),
    (1.00, (244, 248, 250, int(0.48 * 255))),
]


def find_chrome():
    for path in DEFAULT_CHROME_PATHS:
        if Path(path).exists():
            return path
    raise FileNotFoundError("Chrome/Edge executable not found in common install paths.")


def as_file_url(path: Path) -> str:
    return path.resolve().as_uri()


def read_html(path: Path) -> str:
    return path.read_text(encoding="utf-8")


def parse_slides(html: str):
    soup = BeautifulSoup(html, "html.parser")
    slides = []
    for section in soup.select("section.slide"):
        img = section.select_one("img.slide-bg")
        classes = section.get("class", [])
        slides.append(
            {
                "classes": classes,
                "is_light": "light" in classes,
                "bg_src": img.get("src") if img else None,
            }
        )
    return slides


def cover_crop_image(src: Path, out: Path, width: int, height: int):
    img = Image.open(src).convert("RGB")
    ratio = max(width / img.width, height / img.height)
    resized = img.resize((round(img.width * ratio), round(img.height * ratio)), Image.Resampling.LANCZOS)
    left = (resized.width - width) // 2
    top = (resized.height - height) // 2
    cropped = resized.crop((left, top, left + width, top + height))
    cropped = ImageEnhance.Color(cropped).enhance(1.08)
    cropped = ImageEnhance.Contrast(cropped).enhance(1.08)
    cropped.save(out)


def interp(a, b, t):
    return round(a + (b - a) * t)


def gradient_mask(out: Path, width: int, height: int, is_light: bool):
    stops = LIGHT_STOPS if is_light else DARK_STOPS
    img = Image.new("RGBA", (width, height))
    px = img.load()
    for x in range(width):
        pos = x / max(1, width - 1)
        lo, hi = stops[0], stops[-1]
        for i in range(len(stops) - 1):
            if stops[i][0] <= pos <= stops[i + 1][0]:
                lo, hi = stops[i], stops[i + 1]
                break
        span = max(0.00001, hi[0] - lo[0])
        t = (pos - lo[0]) / span
        rgba = tuple(interp(lo[1][c], hi[1][c], t) for c in range(4))
        for y in range(height):
            px[x, y] = rgba
    img.save(out)


def run_chrome_screenshot(chrome: str, html_file: Path, out_png: Path, width: int, height: int):
    cmd = [
        chrome,
        "--headless=new",
        "--disable-gpu",
        "--no-sandbox",
        "--allow-file-access-from-files",
        f"--window-size={width},{height}",
        f"--screenshot={str(out_png)}",
        as_file_url(html_file),
    ]
    subprocess.run(cmd, check=False, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    for _ in range(20):
        if out_png.exists() and out_png.stat().st_size > 0:
            return
        time.sleep(0.25)
    raise RuntimeError(f"Screenshot was not created: {out_png}")


def inject_capture_css(html: str, input_dir: Path, slide_index: int, width: int, height: int, slide_count: int):
    base = f'<base href="{as_file_url(input_dir)}/">'
    extra = f"""
<style id="hybrid-capture-style">
html, body {{
  width:{width}px !important;
  height:{height}px !important;
  margin:0 !important;
  overflow:hidden !important;
  background:#05070a !important;
}}
.deck {{
  width:{width * slide_count}px !important;
  height:{height}px !important;
  transform:translateX(-{width * slide_index}px) !important;
  transition:none !important;
}}
.slide {{
  width:{width}px !important;
  height:{height}px !important;
  flex:0 0 {width}px !important;
}}
.progress {{ display:none !important; }}
</style>
"""
    return html.replace("<head>", f"<head>\n{base}", 1).replace("</head>", extra + "\n</head>", 1)


def export_pdf_playwright(input_html: Path, out_pdf: Path, width_mm: float, height_mm: float, timeout_ms: int):
    try:
        from playwright.sync_api import sync_playwright
    except Exception as exc:
        raise RuntimeError("Playwright is required for high-fidelity PDF output.") from exc

    html_url = as_file_url(input_html)
    with sync_playwright() as pw:
        browser = pw.chromium.launch(
            headless=True,
            args=["--no-sandbox", "--disable-setuid-sandbox", "--disable-gpu", "--font-render-hinting=none"],
        )
        page = browser.new_page(viewport={"width": 1920, "height": 1080}, device_scale_factor=1)
        try:
            page.goto(html_url, wait_until="networkidle", timeout=timeout_ms)
        except Exception:
            page.goto(html_url, wait_until="domcontentloaded", timeout=timeout_ms)
        page.evaluate(
            """async () => {
              const imgs = [...document.querySelectorAll('img')];
              await Promise.allSettled(imgs.map(img => img.complete ? Promise.resolve() : new Promise(r => { img.onload=r; img.onerror=r; })));
              try { await document.fonts.ready; } catch(e) {}
            }"""
        )
        page.wait_for_timeout(1000)
        pdf_data = page.pdf(
            width=f"{width_mm}mm",
            height=f"{height_mm}mm",
            margin={"top": "0mm", "right": "0mm", "bottom": "0mm", "left": "0mm"},
            print_background=True,
            prefer_css_page_size=True,
            scale=1.0,
        )
        browser.close()
    out_pdf.write_bytes(pdf_data)


def extract_text_boxes(input_html: Path, out_json: Path, width: int, height: int, slide_count: int, timeout_ms: int):
    try:
        from playwright.sync_api import sync_playwright
    except Exception as exc:
        raise RuntimeError("Playwright is required for editable text extraction.") from exc

    js = """
    ({width, height}) => {
      function visible(el){
        const cs = getComputedStyle(el);
        const r = el.getBoundingClientRect();
        return cs.display !== 'none' && cs.visibility !== 'hidden' && parseFloat(cs.opacity || '1') > 0.01 && r.width > 2 && r.height > 2;
      }
      function rgbToHex(color){
        const m = color.match(/rgba?\\((\\d+),\\s*(\\d+),\\s*(\\d+)/);
        if(!m) return '000000';
        return [m[1],m[2],m[3]].map(v => Number(v).toString(16).padStart(2,'0')).join('').toUpperCase();
      }
      const slides = [...document.querySelectorAll('section.slide')];
      return slides.map((slide, pageIndex) => {
        const deck = document.querySelector('.deck');
        if(deck) deck.style.transform = `translateX(-${pageIndex * width}px)`;
        document.body.offsetHeight;
        const nodes = [...slide.querySelectorAll('.kicker,.title,.subtitle,.cover-ribbon span,.cover-ribbon b,.slash-panel b,.slash-panel span,.large-note,.editorial-strip strong,.editorial-strip b,.editorial-strip span,.mega b,.mega span,.mega em,.poster-side b,.poster-side span,.data-footnote,.proof-board b,.proof-board span,.pin b,.pin span,.chain-slices em,.chain-slices b,.chain-slices span,.chain-note,.a-list span,.target-logic b,.target-logic span,.chrome span,.foot span,.pager')];
        return nodes.filter(el => visible(el) && (el.innerText || '').trim()).map(el => {
          const r = el.getBoundingClientRect();
          const sr = slide.getBoundingClientRect();
          const cs = getComputedStyle(el);
          const text = (el.innerText || '').replace(/\\s+/g, ' ').trim();
          return {
            text,
            left: r.left - sr.left,
            top: r.top - sr.top,
            width: r.width,
            height: r.height,
            fontSize: parseFloat(cs.fontSize) || 16,
            fontFamily: cs.fontFamily || '',
            fontWeight: cs.fontWeight || '400',
            color: rgbToHex(cs.color || 'rgb(0,0,0)'),
            align: cs.textAlign || 'left',
            italic: cs.fontStyle === 'italic',
            lineHeight: cs.lineHeight
          };
        });
      });
    }
    """
    with sync_playwright() as pw:
        browser = pw.chromium.launch(headless=True, args=["--no-sandbox", "--disable-gpu"])
        page = browser.new_page(viewport={"width": width, "height": height}, device_scale_factor=1)
        try:
            page.goto(as_file_url(input_html), wait_until="networkidle", timeout=timeout_ms)
        except Exception:
            page.goto(as_file_url(input_html), wait_until="domcontentloaded", timeout=timeout_ms)
        page.evaluate(
            """async () => {
              await Promise.allSettled([...document.querySelectorAll('img')].map(img => img.complete ? Promise.resolve() : new Promise(r => { img.onload=r; img.onerror=r; })));
              try { await document.fonts.ready; } catch(e) {}
            }"""
        )
        page.add_style_tag(
            content=f"""
            html,body{{width:{width}px!important;height:{height}px!important;margin:0!important;overflow:hidden!important}}
            .deck{{width:{width * slide_count}px!important;height:{height}px!important;transition:none!important}}
            .slide{{width:{width}px!important;height:{height}px!important;flex:0 0 {width}px!important}}
            """
        )
        data = page.evaluate(js, {"width": width, "height": height})
        browser.close()
    out_json.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
    return data


def rgb(hex_color: str) -> RGBColor:
    value = hex_color.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def choose_font(font_family: str, default: str = "Noto Sans SC") -> str:
    fam = font_family.lower()
    if "serif" in fam or "simsun" in fam or "noto serif" in fam:
        return "Noto Serif SC"
    if "mono" in fam or "consolas" in fam or "jetbrains" in fam:
        return "Consolas"
    return default


def make_editable_pptx(out_pptx: Path, text_boxes, layers_dir: Path, slide_count: int, width: int, height: int):
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank = prs.slide_layouts[6]
    sx = 16 / width
    sy = 9 / height
    for i in range(slide_count):
        slide = prs.slides.add_slide(blank)
        n = i + 1
        slide.shapes.add_picture(str(layers_dir / "background" / f"slide-{n:02d}-background.png"), 0, 0, Inches(16), Inches(9))
        slide.shapes.add_picture(str(layers_dir / "mask" / f"slide-{n:02d}-mask.png"), 0, 0, Inches(16), Inches(9))
        for item in text_boxes[i]:
            text = item["text"]
            if not text:
                continue
            left = Inches(max(0, item["left"] * sx))
            top = Inches(max(0, item["top"] * sy))
            box_w = Inches(max(0.1, item["width"] * sx + 0.05))
            box_h = Inches(max(0.1, item["height"] * sy + 0.06))
            shape = slide.shapes.add_textbox(left, top, box_w, box_h)
            shape.fill.background()
            shape.line.fill.background()
            tf = shape.text_frame
            tf.clear()
            tf.margin_left = 0
            tf.margin_right = 0
            tf.margin_top = 0
            tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.text = text
            if item.get("align") == "center":
                p.alignment = PP_ALIGN.CENTER
            elif item.get("align") == "right":
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT
            for run in p.runs:
                run.font.name = choose_font(item.get("fontFamily", ""))
                run.font.size = Pt(max(5, item["fontSize"] * 0.72))
                run.font.bold = str(item.get("fontWeight", "")).lower() in {"bold", "700", "800", "900"} or (
                    str(item.get("fontWeight", "400")).isdigit() and int(item.get("fontWeight", "400")) >= 650
                )
                run.font.italic = bool(item.get("italic"))
                run.font.color.rgb = rgb(item.get("color", "000000"))
    if len(prs.slides) > slide_count:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]
    prs.save(out_pptx)


def main():
    parser = argparse.ArgumentParser(description="Hybrid HTML deck export: Playwright PDF + layered background/mask + editable PPTX text.")
    parser.add_argument("input_html", type=Path)
    parser.add_argument("-o", "--out-dir", type=Path, required=True)
    parser.add_argument("--width", type=int, default=1600)
    parser.add_argument("--height", type=int, default=900)
    parser.add_argument("--timeout", type=int, default=30000)
    parser.add_argument("--chrome", default=None)
    args = parser.parse_args()

    input_html = args.input_html.resolve()
    out_dir = args.out_dir.resolve()
    out_dir.mkdir(parents=True, exist_ok=True)
    layers_dir = out_dir / "layers"
    for sub in ["background", "mask", "full", "source-background"]:
        (layers_dir / sub).mkdir(parents=True, exist_ok=True)
    tmp_dir = out_dir / "_tmp"
    tmp_dir.mkdir(parents=True, exist_ok=True)

    chrome = args.chrome or find_chrome()
    html = read_html(input_html)
    slides = parse_slides(html)
    if not slides:
        raise RuntimeError("No section.slide elements found.")

    for idx, slide in enumerate(slides):
        n = idx + 1
        bg_path = (input_html.parent / slide["bg_src"]).resolve()
        shutil.copy2(bg_path, layers_dir / "source-background" / f"slide-{n:02d}-source{bg_path.suffix}")
        cover_crop_image(bg_path, layers_dir / "background" / f"slide-{n:02d}-background.png", args.width, args.height)
        gradient_mask(layers_dir / "mask" / f"slide-{n:02d}-mask.png", args.width, args.height, slide["is_light"])
        full_html = tmp_dir / f"capture-full-{n:02d}.html"
        full_html.write_text(inject_capture_css(html, input_html.parent, idx, args.width, args.height, len(slides)), encoding="utf-8")
        run_chrome_screenshot(chrome, full_html, layers_dir / "full" / f"slide-{n:02d}-full.png", args.width, args.height)

    pdf_path = out_dir / "html-deck-16x9-playwright.pdf"
    export_pdf_playwright(input_html, pdf_path, 406.4, 228.6, args.timeout)
    pdf_reader = PdfReader(str(pdf_path))

    text_json = out_dir / "editable-text-boxes.json"
    text_boxes = extract_text_boxes(input_html, text_json, args.width, args.height, len(slides), args.timeout)
    pptx_path = out_dir / "html-deck-16x9-editable-text.pptx"
    make_editable_pptx(pptx_path, text_boxes, layers_dir, len(slides), args.width, args.height)

    report = out_dir / "hybrid-export-report.txt"
    report.write_text(
        "\n".join(
            [
                f"input={input_html}",
                f"slides={len(slides)}",
                f"pdf={pdf_path}",
                f"pdf_pages={len(pdf_reader.pages)}",
                f"pdf_first_page_size={float(pdf_reader.pages[0].mediabox.width)}x{float(pdf_reader.pages[0].mediabox.height)}pt",
                f"pptx={pptx_path}",
                f"pptx_mode=background image + mask image + editable text boxes",
                f"text_box_count={sum(len(x) for x in text_boxes)}",
                f"layer_assets={layers_dir}",
                f"chrome={chrome}",
            ]
        ),
        encoding="utf-8",
    )
    print(report.read_text(encoding="utf-8"))


if __name__ == "__main__":
    try:
        main()
    except Exception as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        raise
